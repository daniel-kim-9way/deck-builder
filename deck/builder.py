# -*- coding: utf-8 -*-
"""
Deck — 덱 빌더. 테마(폰트+스타일)와 덱 유형을 받아 슬라이드를 만든다.
16:9 (1280x720 px). 출력은 <FT_OUTPUT_ROOT 또는 cwd/OUTPUT>/deck-builder/<프로젝트>/.
"""
import os
from pptx import Presentation
import tokens as T
import theme as theme_mod
from helpers import px, rgb, text


def output_root():
    # FT_OUTPUT_ROOT 가 있으면 그걸, 없으면 현재 작업폴더/OUTPUT (이식성 유지)
    return os.environ.get("FT_OUTPUT_ROOT") or os.path.join(os.getcwd(), "OUTPUT")


def project_dir(project, sub=None):
    # 결과물은 <root>/deck-builder/<project>[/<sub>] — 스킬별로 묶는다
    parts = [output_root(), "deck-builder", project] + ([sub] if sub else [])
    d = os.path.join(*parts)
    os.makedirs(d, exist_ok=True)
    return d


class Deck:
    def __init__(self, theme=None, kind=None, brand=""):
        """theme: 프리셋 이름/파일경로/dict (생략 시 기본 cream-coral).
        kind: 'proposal'|'report'|'lecture'|'guidebook' (푸터·계획 가이드용)."""
        if theme is not None:
            theme_mod.use(theme)
        self.kind = kind
        self.brand = brand
        self.prs = Presentation()
        self.prs.slide_width = px(T.SLIDE_W_PX)
        self.prs.slide_height = px(T.SLIDE_H_PX)
        self.blank = self.prs.slide_layouts[6]
        self._page = 0

    def add(self, bg="canvas"):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(bg)
        return s

    def footer(self, slide, dark=False, label=None):
        self._page += 1
        c = "on_dark_soft" if dark else "muted_soft"
        lab = self.brand if label is None else label
        if lab:
            text(slide, 80, 685, 500, 18, lab, type_token="body_sm",
                 color=c, valign="middle", size_px=11)
        text(slide, T.SLIDE_W_PX - 80 - 60, 685, 60, 18, f"{self._page:02d}",
             type_token="body_sm", color=c, align="right", valign="middle",
             size_px=11)

    def save(self, path):
        self.prs.save(path)
        return path

    def save_project(self, project, filename="deck.pptx"):
        """<root>/deck-builder/<project>/<filename> 에 저장 (root=FT_OUTPUT_ROOT 또는 cwd/OUTPUT)."""
        path = os.path.join(project_dir(project), filename)
        self.prs.save(path)
        return path
