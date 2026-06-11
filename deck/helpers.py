# -*- coding: utf-8 -*-
"""
python-pptx 드로잉 헬퍼.
px 좌표계(1280x720)로 작업하고 내부에서 EMU/Pt로 변환한다.
원본 디자인 시스템의 색/타이포/radius 토큰을 충실히 렌더하기 위한 저수준 유틸.
"""
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import tokens as T

EMU_PER_PX = 9525          # 1280px*9525 = 12192000 EMU = 13.333in (16:9 정확 일치)
PT_PER_PX = 0.75           # 1280x720 캔버스에서 px → pt


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def pt(v):
    return Pt(v * PT_PER_PX)


def rgb(key_or_hex):
    h = T.COLORS.get(key_or_hex, key_or_hex)
    return RGBColor.from_string(h)


# ---------------------------------------------------------------------------
# 도형
# ---------------------------------------------------------------------------
def _kill_theme_shadow(sp):
    """테마 style effectRef(그림자) 제거 — 컬러블록 우선 철학."""
    sp.shadow.inherit = False  # 빈 <a:effectLst/> 삽입
    style = sp._element.find(qn('p:style'))
    if style is not None:
        eff = style.find(qn('a:effectRef'))
        if eff is not None:
            eff.set('idx', '0')  # 테마 효과(그림자) 비활성


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         radius=0, shadow=False):
    """둥근/직각 사각형. radius(px)>0 이면 rounded."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    if radius > 0:
        # adjustment = radius / (shorter side) , 0~0.5
        frac = min(0.5, radius / max(1, min(w, h)))
        try:
            sp.adjustments[0] = frac
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    # 다크/hairline 테마: surface_card 카드는 자동으로 hairline 보더 (Voltagent 스타일)
    if line is None and fill == "surface_card" and getattr(T, "CARD_HAIRLINE", False):
        line = "hairline"
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    _kill_theme_shadow(sp)
    if shadow:
        _soft_shadow(sp)
    return sp


def oval(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(w), px(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(line_w)
    _kill_theme_shadow(sp)
    return sp


def _soft_shadow(sp):
    """0 1px 3px rgba(20,20,19,0.08) 근사 — 거의 안 쓰지만 hover 상태용."""
    spPr = sp._element.spPr
    effLst = spPr.find(qn('a:effectLst'))
    if effLst is None:
        effLst = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(effLst)
    sh = effLst.makeelement(qn('a:outerShdw'),
                            {'blurRad': '38100', 'dist': '12700',
                             'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '141413'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '8000'})
    clr.append(alpha); sh.append(clr); effLst.append(sh)


# ---------------------------------------------------------------------------
# 이미지 (스크린샷 / 가이드북)
# ---------------------------------------------------------------------------
def picture(slide, x, y, w, h, path=None, frame=True, radius=10,
            caption="스크린샷"):
    """이미지 박스. path가 있으면 비율 유지 fit + 중앙 배치, 없으면 플레이스홀더.
    가이드북의 '우측 스크린샷' 패턴용. 프레임(hairline) 선택."""
    import os
    bx, by, bw, bh = px(x), px(y), px(w), px(h)
    if path and os.path.isfile(path):
        if frame:
            rect(slide, x, y, w, h, fill="surface_soft", line="hairline",
                 radius=radius)
        pic = slide.shapes.add_picture(path, bx, by)
        nw, nh = pic.width, pic.height
        scale = min(bw / nw, bh / nh)
        pic.width = int(nw * scale)
        pic.height = int(nh * scale)
        pic.left = int(bx + (bw - pic.width) / 2)
        pic.top = int(by + (bh - pic.height) / 2)
        return pic
    # 플레이스홀더 (이미지 미제공 시에도 레이아웃이 깨지지 않게)
    rect(slide, x, y, w, h, fill="surface_card", line="hairline", radius=radius)
    oval(slide, x + w / 2 - 28, y + h / 2 - 44, 56, 56, fill="surface_soft",
         line="hairline")
    text(slide, x + w / 2 - 28, y + h / 2 - 44, 56, 56, "🖼",
         color="muted", align="center", valign="middle", size_px=24)
    text(slide, x, y + h / 2 + 26, w, 28, caption, type_token="caption_upper",
         color="muted_soft", align="center")
    return None


# ---------------------------------------------------------------------------
# 텍스트
# ---------------------------------------------------------------------------
def _apply_run(run, type_token=None, size_px=None, weight=None,
               color="ink", tracking_px=None, mono=False):
    if type_token:
        s, w, lh, tr = T.TYPE[type_token]
        size_px = size_px or s
        weight = weight or w
        tracking_px = tr if tracking_px is None else tracking_px
    size_px = size_px or 16
    weight = weight or "Regular"
    tracking_px = tracking_px or 0

    # 슬라이드 가독성 보정: 작은 텍스트일수록 더 많이 키운다.
    # 28px로 수렴(28px 이상=디스플레이/빅넘버는 그대로). 연속 함수라 역전 없음.
    #   11px → 19.5px(~14.6pt) · 14px → 21px · 16px → 22px · 22px → 25px
    if size_px < 28:
        size_px = size_px + (28 - size_px) * 0.5

    fam, bold = T.WEIGHT_MAP.get(weight, (T.FONT_BODY, False))
    # 디스플레이 토큰이고 display 폰트가 본문과 다르면 display 폰트 사용
    if (type_token and type_token.startswith("display")
            and T.FONT_DISPLAY != T.FONT_BODY and fam == T.FONT_BODY):
        fam = T.FONT_DISPLAY
        bold = weight in T.HEAVY
    if mono:
        fam = T.FONT_MONO; bold = False
    run.font.name = fam
    run.font.size = pt(size_px)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    # 한글/라틴 모두 동일 폰트 지정 (ea/cs)
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', fam)
    # 자간 (spc, 1/100 pt). px → pt → *100
    spc = int(round(tracking_px * PT_PER_PX * 100))
    if spc != 0:
        rPr.set('spc', str(spc))


def text(slide, x, y, w, h, content, type_token="body_md", color="ink",
         align="left", valign="top", size_px=None, weight=None,
         tracking_px=None, line_spacing=None, mono=False, wrap=True):
    """단일 스타일 텍스트박스."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[valign]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    # line height
    if line_spacing is None and type_token:
        line_spacing = T.TYPE[type_token][2]
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run(); run.text = content
    _apply_run(run, type_token, size_px, weight, color, tracking_px, mono)
    return tb


def bullets(slide, x, y, w, h, items, type_token="body_md", color="body",
            gap_px=6, marker="•", marker_color="primary"):
    """리스트. items: [str, ...] 또는 [(marker_char, str)]."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    s, w_, lh, tr = T.TYPE[type_token]
    for i, it in enumerate(items):
        mk, txt = (marker, it) if isinstance(it, str) else it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = lh
        p.space_after = pt(gap_px)
        if mk:
            rm = p.add_run(); rm.text = mk + "  "
            _apply_run(rm, type_token, color=marker_color)
        rt = p.add_run(); rt.text = txt
        _apply_run(rt, type_token, color=color)
    return tb


def line_div(slide, x, y, w, color="hairline", weight=1.0, dash=None):
    """수평 헤어라인. dash='dash'|'dot' 가능."""
    ln = slide.shapes.add_connector(2, px(x), px(y), px(x + w), px(y))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    if dash:
        d = ln.line._get_or_add_ln()
        pd = d.makeelement(qn('a:prstDash'), {'val': dash})
        d.append(pd)
    ln.shadow.inherit = False
    return ln


def diag(slide, x1, y1, x2, y2, color="muted_soft", weight=1.0, dash="dash"):
    """대각선 커넥터 (주석 화살표용)."""
    ln = slide.shapes.add_connector(2, px(x1), px(y1), px(x2), px(y2))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    if dash:
        d = ln.line._get_or_add_ln()
        d.append(d.makeelement(qn('a:prstDash'), {'val': dash}))
    ln.shadow.inherit = False
    return ln


def vline(slide, x, y, h, color="hairline", weight=1.0, dash=None):
    """수직 라인."""
    ln = slide.shapes.add_connector(2, px(x), px(y), px(x), px(y + h))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    if dash:
        d = ln.line._get_or_add_ln()
        d.append(d.makeelement(qn('a:prstDash'), {'val': dash}))
    ln.shadow.inherit = False
    return ln


# ---------------------------------------------------------------------------
# 브랜드 마크 (Anthropic 4-스포크 라디얼 글리프 근사)
# ---------------------------------------------------------------------------
def spike_mark(slide, cx, cy, size, color="ink"):
    """중심(cx,cy) 기준 size(px) 별/방사형 마크. 4개 막대를 45도씩 회전."""
    bar_w = size * 0.14
    for deg in (0, 45, 90, 135):
        sp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            px(cx - bar_w / 2), px(cy - size / 2), px(bar_w), px(size))
        try:
            sp.adjustments[0] = 0.5
        except Exception:
            pass
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
        sp.line.fill.background(); _kill_theme_shadow(sp)
        sp.rotation = deg
    return None


def wordmark(slide, x, y, label="Claude", mark_color="ink", text_color="ink",
             size_px=22):
    """스파이크 마크 + 워드마크."""
    spike_mark(slide, x + size_px * 0.5, y + size_px * 0.5, size_px * 0.9,
               color=mark_color)
    text(slide, x + size_px * 1.25, y, 300, size_px * 1.4, label,
         type_token="title_lg", color=text_color, valign="middle",
         size_px=size_px, weight="SemiBold")


# ---------------------------------------------------------------------------
# 네이티브 표 (편집 가능)
# ---------------------------------------------------------------------------
def _cell_border(cell, color="hairline", w_pt=1.0,
                 edges=("L", "R", "T", "B")):
    """셀 테두리 — tcPr 스키마 순서(lnL,lnR,lnT,lnB가 맨 앞) 준수."""
    tcPr = cell._tc.get_or_add_tcPr()
    val = T.COLORS.get(color, color)
    # 기존 ln 제거
    for e in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tcPr.findall(qn(e)):
            tcPr.remove(el)
    # 역순 삽입 → 최종 L,R,T,B 순서로 맨 앞 배치
    order = [e for e in ("L", "R", "T", "B") if e in edges]
    for edge in reversed(order):
        ln = tcPr.makeelement(qn("a:ln" + edge),
                              {"w": str(int(w_pt * 12700)), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": val})
        fill.append(clr)
        ln.append(fill)
        tcPr.insert(0, ln)


def add_table(slide, x, y, w, h, data, col_widths=None,
              header_fill="surface_dark", header_color="on_dark",
              highlight_col=None, highlight_fill="surface_cream_strong",
              body_fill="canvas", body_color="body", border="hairline",
              header_size=14, body_size=14, row_h=None,
              label_col=False, label_fill="surface_card"):
    """
    data: 2차원 리스트 (행 x 열). data[0]가 헤더.
    highlight_col: 강조할 열 인덱스(0-base) — 9WAY 컬럼 등.
    """
    rows = len(data)
    cols = len(data[0])
    gf = slide.shapes.add_table(rows, cols, px(x), px(y), px(w), px(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 열 폭
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = px(cw)
    # 행 높이
    rh = row_h or (h / rows)
    for ri in range(rows):
        tbl.rows[ri].height = px(rh)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = px(12)
            cell.margin_right = px(12)
            cell.margin_top = px(6)
            cell.margin_bottom = px(6)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 배경
            if ri == 0:
                fill_c, txt_c, weight = header_fill, header_color, "SemiBold"
            elif highlight_col is not None and ci == highlight_col:
                fill_c, txt_c, weight = highlight_fill, "ink", "SemiBold"
            elif label_col and ci == 0:
                fill_c, txt_c, weight = label_fill, "ink", "SemiBold"
            else:
                fill_c, txt_c, weight = body_fill, body_color, "Regular"
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill_c)
            # 텍스트
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 or ri == 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            _apply_run(r, size_px=(header_size if ri == 0 else body_size),
                       weight=weight, color=txt_c)
            _cell_border(cell, color=border, w_pt=1.0)
    return tbl
